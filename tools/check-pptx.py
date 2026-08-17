#!/usr/bin/env python3
"""Structural check on a generated status board, using a parser that is not ours.

    pip install python-pptx
    python3 tools/check-pptx.py /tmp/trip/board.pptx

Opening the file through python-pptx exercises the content types and relationships;
a package PowerPoint would reject generally fails here first. Then it checks the
slide count, that nothing lands off the canvas, and that slide 2 carries a table.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

path = sys.argv[1] if len(sys.argv) > 1 else 'board.pptx'
p = Presentation(path)
fails = []

W, H = p.slide_width / 914400, p.slide_height / 914400
if (round(W, 2), round(H, 2)) != (13.33, 7.5):
    fails.append(f"slide size is {W:.2f}x{H:.2f}, expected 13.33x7.5")

slides = list(p.slides)
if len(slides) != 2:
    fails.append(f"{len(slides)} slides, expected 2")

for i, s in enumerate(slides, 1):
    for sh in s.shapes:
        x, y = Emu(sh.left).inches, Emu(sh.top).inches
        r, b = x + Emu(sh.width).inches, y + Emu(sh.height).inches
        if x < -0.01 or y < -0.01 or r > W + 0.01 or b > H + 0.01:
            fails.append(f"slide {i}: shape off the canvas at {x:.2f},{y:.2f} -> {r:.2f},{b:.2f}")

if slides:
    texts = [sh.text_frame.text for sh in slides[0].shapes if sh.has_text_frame]
    if 'NODE STATUS BOARD' not in texts:
        fails.append("slide 1 has no title")
    if not any(sh.has_table for sh in slides[1].shapes):
        fails.append("slide 2 has no table")

print('\n'.join('FAIL  ' + f for f in fails) if fails else f"PASS  {path} opens clean, 2 slides, nothing off-canvas")
sys.exit(1 if fails else 0)
